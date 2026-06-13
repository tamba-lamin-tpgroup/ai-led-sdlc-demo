"""
Generic Audit PowerPoint Generator.

Reads CSV and markdown outputs from /seo-audit, /crawl-site, and
/tech-stack-audit and produces a professional, branded PowerPoint deck.

Always includes a Methodology slide. Fully configurable via CLI args.

Usage:
    python3 generate_audit_pptx.py --input ./seo-audit-output
    python3 generate_audit_pptx.py --input ./seo-audit-output --client "Acme" --company "Agency"
    python3 generate_audit_pptx.py --input ./seo-audit-output --brand-color "#E63946"
"""

import argparse
import csv
import json
import os
import sys
from collections import Counter
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import pptx.oxml.ns as nsmap
    from lxml import etree
except ImportError:
    sys.exit("python-pptx is required: pip3 install python-pptx --break-system-packages -q")

# ---------------------------------------------------------------------------
# Slide dimensions (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Arg parser
# ---------------------------------------------------------------------------

def build_parser():
    p = argparse.ArgumentParser(description="Audit PowerPoint Generator",
                                formatter_class=argparse.ArgumentDefaultsHelpFormatter)
    p.add_argument("--input", required=True, help="Audit output directory")
    p.add_argument("--output", default=None, help="Output .pptx path")
    p.add_argument("--client", default="Website", help="Client/project name")
    p.add_argument("--company", default="SEO Audit", help="Auditing company name")
    p.add_argument("--brand-color", default="#1a1a2e", help="Primary brand hex color")
    p.add_argument("--accent-color", default=None, help="Accent hex color")
    p.add_argument("--logo", default=None, help="Logo image file path")
    p.add_argument("--url", default=None, help="Target URL override")
    p.add_argument("--date", default=None, help="Audit date string override")
    p.add_argument("--no-geo", action="store_true", help="Skip GEO slides")
    p.add_argument("--confidential", action="store_true", help="Add confidentiality footer")
    return p


# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def lighten(hex_str: str, factor: float = 0.85) -> RGBColor:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


# Severity colors (fixed regardless of brand)
SEV_COLORS = {
    "Critical": RGBColor(0xDA, 0x29, 0x1C),
    "High":     RGBColor(0xFF, 0xC0, 0x00),
    "Medium":   RGBColor(0x00, 0x70, 0xC0),
    "Low":      RGBColor(0xA5, 0xA5, 0xA5),
}
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x2C, 0x2C)
GREEN = RGBColor(0x00, 0xB0, 0x50)
LIGHT_GREY = RGBColor(0xF4, 0xF4, 0xF4)


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_slide_header(slide, title: str, brand_rgb: RGBColor,
                     subtitle: str = "", company: str = "") -> None:
    """Dark header bar across the top of every content slide."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), brand_rgb)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.15), Inches(11.5), Inches(0.7),
                 font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.75), Inches(8), Inches(0.4),
                     font_size=12, color=WHITE)
    if company:
        add_text_box(slide, company,
                     Inches(10), Inches(0.82), Inches(3), Inches(0.3),
                     font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide, company: str, confidential: bool,
               brand_rgb: RGBColor) -> None:
    """Thin footer bar at the bottom of every slide."""
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), brand_rgb)
    txt = f"CONFIDENTIAL — {company}" if confidential else company
    add_text_box(slide, txt,
                 Inches(0.2), Inches(7.22), Inches(13), Inches(0.25),
                 font_size=8, color=WHITE)


def progress_bar(slide, label: str, value: int, max_value: int,
                 left, top, width, height,
                 fill_color: RGBColor, bg_color: RGBColor = LIGHT_GREY) -> None:
    """Horizontal progress bar with percentage label."""
    add_rect(slide, left, top, width, height, bg_color)
    pct = min(1.0, value / max_value) if max_value else 0
    if pct > 0:
        add_rect(slide, left, top, int(width * pct), height, fill_color)
    pct_text = f"{int(pct * 100)}%"
    add_text_box(slide, f"{label}  {pct_text}",
                 left, top - Inches(0.25), width, Inches(0.25),
                 font_size=11, color=DARK)


# ---------------------------------------------------------------------------
# Data loaders
# ---------------------------------------------------------------------------

def load_csv(path: str) -> list[dict]:
    if not os.path.exists(path):
        return []
    with open(path, encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def csv_kv(path: str) -> dict[str, str]:
    rows = load_csv(path)
    return {r.get("metric", r.get("signal", "")): str(r.get("value", "")) for r in rows}


def load_audit_data(input_dir: str) -> dict:
    csv_dir = os.path.join(input_dir, "csv")
    rep_dir = os.path.join(input_dir, "reports")

    seo = load_csv(os.path.join(csv_dir, "seo_crawl_raw.csv"))
    recs = load_csv(os.path.join(csv_dir, "recommendations.csv"))
    budget = csv_kv(os.path.join(csv_dir, "crawl_budget_analysis.csv"))
    find = load_csv(os.path.join(csv_dir, "findability_analysis.csv"))
    find_kv = {r.get("metric", ""): r.get("value", "") for r in find}
    geo = load_csv(os.path.join(csv_dir, "geo_structured_data.csv"))
    geo_entity = load_csv(os.path.join(csv_dir, "geo_entity_relationships.csv"))
    geo_ai = {r.get("signal", ""): r.get("value", "")
              for r in load_csv(os.path.join(csv_dir, "geo_ai_crawlability.csv"))}
    errors = load_csv(os.path.join(csv_dir, "seo_errors.csv"))
    dup_titles = load_csv(os.path.join(csv_dir, "seo_duplicate_titles.csv"))

    total = len(seo)
    statuses = Counter(str(r.get("http_status", "")) for r in seo)
    total_200 = statuses.get("200", 0)

    def count(field, val="True"):
        return sum(1 for r in seo if str(r.get(field, "")) == val)

    # Try to extract URL from first 200 row
    base_url = ""
    for r in seo:
        if str(r.get("http_status")) == "200" and r.get("url"):
            from urllib.parse import urlparse
            p = urlparse(r["url"])
            base_url = f"{p.scheme}://{p.netloc}"
            break

    return dict(
        total=total,
        total_200=total_200,
        statuses=dict(statuses),
        base_url=base_url,
        missing_title=count("title_missing"),
        missing_desc=count("meta_description_missing"),
        missing_h1=count("h1_missing"),
        missing_canon=count("canonical_missing"),
        canon_conflict=count("canonical_conflict"),
        dup_titles=count("title_duplicate_flag"),
        dup_descs=count("meta_description_duplicate_flag"),
        noindex=count("has_noindex"),
        errors_4xx=sum(v for k, v in statuses.items() if k.startswith("4")),
        errors_5xx=sum(v for k, v in statuses.items() if k.startswith("5")),
        img_alt_pages=sum(1 for r in seo if int(r.get("images_missing_alt", 0) or 0) > 0),
        total_imgs_no_alt=sum(int(r.get("images_missing_alt", 0) or 0) for r in seo),
        recs=recs[:15],
        budget=budget,
        find_kv=find_kv,
        geo_rows=geo,
        geo_entity=geo_entity,
        geo_ai=geo_ai,
        errors=errors,
        dup_title_rows=dup_titles[:5],
        breadcrumb_count=sum(1 for r in seo if "BreadcrumbList" in (r.get("schema_types_present") or "")),
        has_geo_data=len(geo) > 0,
        llms_present=geo_ai.get("llms_txt_present", "False").lower() == "true",
        blocked_ai_bots=geo_ai.get("robots_txt_blocked_ai_bots", "Not checked"),
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_cover(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Full dark background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, brand_rgb)
    # White accent bar (left edge)
    add_rect(slide, Inches(0), Inches(2), Inches(0.08), Inches(3.5), WHITE)
    # Title
    add_text_box(slide, f"{args.client}\nSEO & Web Audit Report",
                 Inches(0.5), Inches(1.8), Inches(10), Inches(2),
                 font_size=36, bold=True, color=WHITE)
    # Subtitle details
    detail = f"{d['base_url'] or args.url or ''}\nAudit Date: {audit_date}\nPrepared by: {args.company}"
    add_text_box(slide, detail,
                 Inches(0.5), Inches(4.2), Inches(10), Inches(1.5),
                 font_size=14, color=lighten(args.brand_color, 0.6))
    # KPI chips at bottom
    chips = [
        (f"{d['total']:,}", "Pages Audited"),
        (f"{d['total_200']:,}", "HTTP 200"),
        (f"{d['errors_4xx'] + d['errors_5xx']:,}", "Errors Found"),
        (f"{len(d['recs'])}", "Recommendations"),
    ]
    chip_w = Inches(2.8)
    for i, (val, label) in enumerate(chips):
        lft = Inches(0.5) + i * chip_w
        add_rect(slide, lft, Inches(5.9), chip_w - Inches(0.1), Inches(1.2),
                 lighten(args.brand_color, 0.15))
        add_text_box(slide, val, lft + Inches(0.1), Inches(5.95),
                     chip_w - Inches(0.2), Inches(0.55),
                     font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, lft + Inches(0.1), Inches(6.48),
                     chip_w - Inches(0.2), Inches(0.3),
                     font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if args.logo and os.path.exists(args.logo):
        slide.shapes.add_picture(args.logo, Inches(10.5), Inches(0.3),
                                 width=Inches(2.5))


def build_agenda(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Agenda", brand_rgb, company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    sections = [
        "1. Methodology",
        "2. Audit Scope & Summary",
        "3. SEO Health Dashboard",
        "4. Top Critical Issues",
        "5. On-Page SEO Findings",
        "6. Technical SEO Findings",
        "7. Image & Accessibility",
    ]
    if d["has_geo_data"] and not args.no_geo:
        sections += ["8. GEO / Structured Data", "9. AI Crawlability"]
    sections += [
        f"{len(sections)  + 1}. Crawl Budget Analysis",
        f"{len(sections)  + 2}. Findability Analysis",
        f"{len(sections)  + 3}. Recommendations & Priority Matrix",
        f"{len(sections)  + 4}. Next Steps",
    ]

    col_size = (len(sections) + 1) // 2
    for i, item in enumerate(sections):
        col = i // col_size
        row = i % col_size
        add_text_box(slide, item,
                     Inches(0.5 + col * 6.5), Inches(1.4 + row * 0.6),
                     Inches(6), Inches(0.55),
                     font_size=14, color=DARK)


def build_methodology(prs, d, brand_rgb, args, audit_date):
    """
    Methodology slide — ALWAYS included. Documents audit approach, phases,
    tools, scope, and limitations.
    """
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Methodology", brand_rgb,
                     subtitle="How this audit was conducted",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    # Left column: phases
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(5.9), Inches(5.8), LIGHT_GREY)
    add_text_box(slide, "Audit Phases",
                 Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.4),
                 font_size=13, bold=True, color=DARK)

    phases = [
        ("Phase 1", "URL Discovery",
         f"Sitemap parsing (index + paginated), robots.txt analysis, BFS link-following fallback. "
         f"Exclusion filters applied. Final crawl list: {d['total']:,} URLs."),
        ("Phase 2", "SEO Crawl",
         f"Full per-URL metadata extraction: title, meta description, canonical, robots directives, "
         f"H1–H3 headings, image alt text, internal/external links, OpenGraph, Twitter Card, "
         f"Schema.org types, HTTP status, response time, page size."),
        ("Phase 3", "GEO / Structured Data Audit",
         "JSON-LD extraction, entity type coverage analysis (Organization, Person, Article, FAQPage, "
         "BreadcrumbList), E-E-A-T signal detection, FAQ schema gap analysis, "
         "robots.txt AI bot rules, llms.txt check." if d["has_geo_data"] and not args.no_geo
         else "Not executed for this audit."),
        ("Phase 4", "Analysis & Reporting",
         "Duplicate title/description detection, crawl budget analysis, findability assessment "
         "(internal link sampling, orphan detection), prioritized recommendations."),
    ]

    y = Inches(1.8)
    for phase_num, phase_name, phase_desc in phases:
        add_rect(slide, Inches(0.5), y, Inches(0.7), Inches(0.28), brand_rgb)
        add_text_box(slide, phase_num, Inches(0.5), y, Inches(0.7), Inches(0.28),
                     font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, phase_name,
                     Inches(1.3), y, Inches(5), Inches(0.28),
                     font_size=11, bold=True, color=DARK)
        add_text_box(slide, phase_desc,
                     Inches(1.3), y + Inches(0.3), Inches(4.8), Inches(0.55),
                     font_size=9, color=DARK)
        y += Inches(1.0)

    # Right column: tools + scope + limitations
    add_rect(slide, Inches(6.6), Inches(1.25), Inches(6.3), Inches(2.6), LIGHT_GREY)
    add_text_box(slide, "Tools & Technology",
                 Inches(6.75), Inches(1.35), Inches(6), Inches(0.35),
                 font_size=13, bold=True, color=DARK)
    tools = [
        "Python 3 — crawler, parser, analysis engine",
        "requests + BeautifulSoup + lxml — HTTP fetch and HTML parsing",
        "extruct — JSON-LD, Microdata, OpenGraph extraction",
        "tqdm — progress monitoring",
        "python-pptx — this report",
    ]
    for i, t in enumerate(tools):
        add_text_box(slide, f"• {t}",
                     Inches(6.75), Inches(1.75) + i * Inches(0.3),
                     Inches(6), Inches(0.28), font_size=9, color=DARK)

    add_rect(slide, Inches(6.6), Inches(3.95), Inches(6.3), Inches(3.1), LIGHT_GREY)
    add_text_box(slide, "Scope & Limitations",
                 Inches(6.75), Inches(4.05), Inches(6), Inches(0.35),
                 font_size=13, bold=True, color=DARK)

    target_label = args.url or d.get("base_url", "Not specified")
    scope_notes = [
        f"Target: {target_label}",
        f"URLs crawled: {d['total']:,} (from sitemap + BFS discovery)",
        f"Audit date: {audit_date}",
        "This audit used a local/development instance (DDEV)" if "ddev" in target_label.lower()
        else "This audit crawled the live production site",
        "SSL verification disabled for local instances",
        "JavaScript-rendered content audited as raw HTML only",
        "WAF/bot blocking on live site not testable from local instance",
        "Screenshots not included (use --screenshots with /crawl-site)",
    ]
    for i, note in enumerate(scope_notes[:7]):
        add_text_box(slide, f"• {note}",
                     Inches(6.75), Inches(4.45) + i * Inches(0.3),
                     Inches(6), Inches(0.28), font_size=9, color=DARK)


def build_scope_summary(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Audit Scope & Summary", brand_rgb, company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    kpis = [
        ("Total URLs", f"{d['total']:,}"),
        ("HTTP 200", f"{d['total_200']:,}"),
        ("4xx Errors", f"{d['errors_4xx']:,}"),
        ("5xx Errors", f"{d['errors_5xx']:,}"),
        ("Redirects Found", str(d.get("redirect_count", "—"))),
        ("Duplicate Titles", f"{d['dup_titles']:,}"),
    ]
    kpi_w = Inches(2.0)
    for i, (label, val) in enumerate(kpis):
        col = i % 3
        row = i // 3
        lft = Inches(0.4) + col * Inches(4.3)
        top = Inches(1.4) + row * Inches(2.0)
        add_rect(slide, lft, top, Inches(4.0), Inches(1.7), LIGHT_GREY)
        add_text_box(slide, val, lft + Inches(0.15), top + Inches(0.2),
                     Inches(3.7), Inches(0.9),
                     font_size=40, bold=True, color=brand_rgb, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, lft + Inches(0.15), top + Inches(1.2),
                     Inches(3.7), Inches(0.4),
                     font_size=13, color=DARK, align=PP_ALIGN.CENTER)

    # HTTP status pie-like summary (text-based)
    status_lines = ", ".join(f"HTTP {k}: {v:,}" for k, v in sorted(d["statuses"].items())
                             if v > 0 and k not in ("", "TIMEOUT", "ERROR", "CONN_ERROR"))
    add_text_box(slide, f"Status distribution: {status_lines}",
                 Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.5),
                 font_size=11, color=DARK)


def build_seo_dashboard(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "SEO Health Dashboard", brand_rgb,
                     subtitle="0 = worst, 100 = best",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    t = d["total_200"] or 1

    def sev_color(bad_rate):
        if bad_rate > 0.5:
            return SEV_COLORS["Critical"]
        if bad_rate > 0.2:
            return SEV_COLORS["High"]
        if bad_rate > 0.05:
            return SEV_COLORS["Medium"]
        return GREEN

    dims = [
        ("Title Tags",          d["missing_title"] + d["dup_titles"] // 2,  t, "Missing + duplicates"),
        ("Meta Descriptions",   d["missing_desc"] + d["dup_descs"] // 2,    t, "Missing + duplicates"),
        ("H1 Headings",         d["missing_h1"],                             t, "Missing H1"),
        ("HTTP Errors",         d["errors_4xx"] + d["errors_5xx"],           t, "4xx + 5xx"),
        ("Canonical Tags",      d["missing_canon"] + d["canon_conflict"],    t, "Missing + conflicts"),
        ("Image Alt Text",      d["img_alt_pages"],                          t, "Pages with missing alt"),
        ("Crawl Budget",        0 if d["errors_4xx"] < 200 else d["errors_4xx"],  max(d["total"], 1), "Error rate"),
    ]

    bar_w = Inches(8)
    for i, (label, bad, total_, hint) in enumerate(dims):
        top = Inches(1.35 + i * 0.75)
        good = max(0, total_ - bad)
        pct = int(100 * good / total_) if total_ else 100
        col = sev_color(bad / total_ if total_ else 0)
        add_text_box(slide, label,
                     Inches(0.4), top + Inches(0.08), Inches(2.8), Inches(0.5),
                     font_size=12, color=DARK)
        # Background bar
        add_rect(slide, Inches(3.3), top + Inches(0.1), bar_w, Inches(0.38), LIGHT_GREY)
        # Fill bar
        if pct > 0:
            add_rect(slide, Inches(3.3), top + Inches(0.1),
                     int(bar_w * pct / 100), Inches(0.38), GREEN if pct > 70 else col)
        add_text_box(slide, f"{pct}%",
                     Inches(11.5), top + Inches(0.1), Inches(1.2), Inches(0.38),
                     font_size=11, bold=True, color=DARK)
        add_text_box(slide, f"{bad:,} affected — {hint}",
                     Inches(3.3), top + Inches(0.5), bar_w, Inches(0.25),
                     font_size=9, color=DARK)


def build_critical_issues(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Top Critical Issues", brand_rgb,
                     subtitle="Ranked by severity and business impact",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    headers = ["#", "Issue", "Severity", "Affected URLs", "Effort", "Impact"]
    col_widths = [Inches(0.4), Inches(5.6), Inches(1.1), Inches(1.4), Inches(0.8), Inches(0.8)]
    col_x = [Inches(0.3)]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)

    # Header row
    header_top = Inches(1.25)
    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        add_rect(slide, cx, header_top, cw, Inches(0.35), brand_rgb)
        add_text_box(slide, h, cx + Inches(0.05), header_top + Inches(0.04),
                     cw - Inches(0.05), Inches(0.3),
                     font_size=10, bold=True, color=WHITE)

    top_recs = [r for r in d["recs"]
                if r.get("severity") in ("Critical", "High")][:10]
    for i, rec in enumerate(top_recs):
        row_top = Inches(1.6 + i * 0.48)
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.3), row_top, Inches(10.2), Inches(0.42), bg)

        sev = rec.get("severity", "Low")
        sev_color = SEV_COLORS.get(sev, DARK)
        row_vals = [
            str(rec.get("priority_rank", i + 1)),
            str(rec.get("issue_title", ""))[:80],
            sev,
            f"{int(rec.get('affected_url_count', 0) or 0):,}",
            str(rec.get("effort", "")),
            str(rec.get("impact", "")),
        ]
        for j, (val, cx, cw) in enumerate(zip(row_vals, col_x, col_widths)):
            text_color = sev_color if j == 2 else DARK
            bold = j == 2
            add_text_box(slide, val,
                         cx + Inches(0.05), row_top + Inches(0.04),
                         cw - Inches(0.1), Inches(0.38),
                         font_size=9, bold=bold, color=text_color)


def build_onpage_seo(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "On-Page SEO Findings", brand_rgb, company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    t = d["total_200"] or 1

    def pct(n): return f"{100*n//t}%"

    metrics = [
        ("Missing Title Tags",          d["missing_title"],  pct(d["missing_title"]),  "Critical"),
        ("Duplicate Title Tags",         d["dup_titles"],     pct(d["dup_titles"]),     "Critical"),
        ("Missing Meta Descriptions",    d["missing_desc"],   pct(d["missing_desc"]),   "Critical"),
        ("Duplicate Meta Descriptions",  d["dup_descs"],      pct(d["dup_descs"]),      "High"),
        ("Missing H1 Headings",          d["missing_h1"],     pct(d["missing_h1"]),     "High"),
        ("Missing Canonical Tags",       d["missing_canon"],  pct(d["missing_canon"]),  "High"),
        ("Canonical Conflicts",          d["canon_conflict"], pct(d["canon_conflict"]), "High"),
        ("Pages with noindex",           d["noindex"],        pct(d["noindex"]),        "Medium"),
    ]

    # Left: metric table
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(7.0), Inches(5.8), LIGHT_GREY)
    headers = ["Finding", "Count", "%", "Severity"]
    col_x = [Inches(0.4), Inches(5.0), Inches(5.9), Inches(6.5)]
    col_w = [Inches(4.5), Inches(0.85), Inches(0.55), Inches(1.0)]

    add_rect(slide, Inches(0.3), Inches(1.25), Inches(7.0), Inches(0.38), brand_rgb)
    for h, cx, cw in zip(headers, col_x, col_w):
        add_text_box(slide, h, cx, Inches(1.28), cw, Inches(0.33),
                     font_size=10, bold=True, color=WHITE)

    for i, (label, count, p, sev) in enumerate(metrics):
        top = Inches(1.65 + i * 0.55)
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.3), top, Inches(7.0), Inches(0.48), bg)
        add_text_box(slide, label, Inches(0.4), top + Inches(0.08),
                     Inches(4.5), Inches(0.38), font_size=10, color=DARK)
        add_text_box(slide, f"{count:,}", Inches(5.0), top + Inches(0.08),
                     Inches(0.85), Inches(0.38), font_size=10, color=DARK)
        add_text_box(slide, p, Inches(5.9), top + Inches(0.08),
                     Inches(0.55), Inches(0.38), font_size=10, color=DARK)
        sev_col = SEV_COLORS.get(sev, DARK)
        add_rect(slide, Inches(6.5), top + Inches(0.1), Inches(0.75), Inches(0.28), sev_col)
        add_text_box(slide, sev, Inches(6.5), top + Inches(0.1),
                     Inches(0.75), Inches(0.28),
                     font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right: top duplicate titles
    add_rect(slide, Inches(7.7), Inches(1.25), Inches(5.2), Inches(5.8), LIGHT_GREY)
    add_text_box(slide, "Top Duplicate Titles",
                 Inches(7.85), Inches(1.35), Inches(4.9), Inches(0.35),
                 font_size=13, bold=True, color=DARK)
    for i, row in enumerate(d.get("dup_title_rows", [])[:5]):
        top = Inches(1.78 + i * 0.9)
        title = str(row.get("title_tag", ""))[:55]
        count = str(row.get("affected_url_count", ""))
        add_rect(slide, Inches(7.85), top, Inches(4.9), Inches(0.78), WHITE)
        add_text_box(slide, f'"{title}"',
                     Inches(7.95), top + Inches(0.05), Inches(4.7), Inches(0.4),
                     font_size=9, color=DARK)
        add_text_box(slide, f"{count} pages",
                     Inches(7.95), top + Inches(0.45), Inches(2), Inches(0.28),
                     font_size=9, bold=True, color=SEV_COLORS["Critical"])


def build_technical_seo(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Technical SEO Findings", brand_rgb, company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    sections = [
        ("HTTP Errors", [
            ("4xx Client Errors", d["errors_4xx"], "High"),
            ("5xx Server Errors", d["errors_5xx"], "Critical"),
        ]),
        ("Redirects", [
            ("Total Redirects", int(d.get("redirect_count", 0)), "Medium"),
            ("Redirect Chains > 2 hops", int(d.get("chain_issues", 0)), "Medium"),
        ]),
        ("Crawl Budget", [
            ("robots.txt has admin rules", None, "Info"),
            ("robots.txt has API endpoint rules", None, "Info"),
        ]),
    ]

    y = Inches(1.35)
    for sec_title, items in sections:
        add_rect(slide, Inches(0.4), y, Inches(12.4), Inches(0.38), brand_rgb)
        add_text_box(slide, sec_title, Inches(0.5), y + Inches(0.04),
                     Inches(12), Inches(0.3),
                     font_size=12, bold=True, color=WHITE)
        y += Inches(0.4)
        for label, count, sev in items:
            add_rect(slide, Inches(0.4), y, Inches(12.4), Inches(0.48), LIGHT_GREY)
            add_text_box(slide, label, Inches(0.55), y + Inches(0.08),
                         Inches(8), Inches(0.38), font_size=11, color=DARK)
            if count is not None:
                add_text_box(slide, f"{count:,}",
                             Inches(9.5), y + Inches(0.08), Inches(1.5), Inches(0.38),
                             font_size=11, bold=True, color=SEV_COLORS.get(sev, DARK))
            sev_col = SEV_COLORS.get(sev, DARK)
            add_rect(slide, Inches(11.2), y + Inches(0.1), Inches(1.2), Inches(0.28), sev_col)
            add_text_box(slide, sev, Inches(11.2), y + Inches(0.1), Inches(1.2), Inches(0.28),
                         font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            y += Inches(0.5)
        y += Inches(0.2)

    # Sample 404 URLs
    error_404 = [r["url"] for r in d.get("errors", [])
                 if str(r.get("http_status")) == "404"][:5]
    if error_404:
        add_text_box(slide, "Sample 404 URLs:",
                     Inches(0.4), y + Inches(0.1), Inches(12), Inches(0.3),
                     font_size=11, bold=True, color=DARK)
        for i, u in enumerate(error_404):
            add_text_box(slide, f"• {u}",
                         Inches(0.6), y + Inches(0.45 + i * 0.35), Inches(12), Inches(0.3),
                         font_size=9, color=SEV_COLORS["High"])


def build_geo_slide(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "GEO / Structured Data Coverage", brand_rgb,
                     subtitle="Schema.org entity coverage across sampled pages",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    geo_entity = d.get("geo_entity", [])
    if not geo_entity:
        add_text_box(slide, "GEO data not available. Run /seo-audit --geo to collect structured data metrics.",
                     Inches(1), Inches(3), Inches(11), Inches(1), font_size=14, color=DARK)
        return

    gt = int(geo_entity[0].get("sample_total", 1)) if geo_entity else 1
    add_text_box(slide, f"Sample: {gt:,} pages",
                 Inches(0.4), Inches(1.25), Inches(5), Inches(0.3),
                 font_size=11, color=DARK)

    headers = ["Entity Type", "Pages With Schema", "Coverage", "Status"]
    col_x = [Inches(0.4), Inches(5.8), Inches(7.6), Inches(9.0)]
    col_w = [Inches(5.3), Inches(1.7), Inches(1.3), Inches(3.8)]

    add_rect(slide, Inches(0.35), Inches(1.55), Inches(12.5), Inches(0.38), brand_rgb)
    for h, cx, cw in zip(headers, col_x, col_w):
        add_text_box(slide, h, cx, Inches(1.58), cw, Inches(0.33),
                     font_size=10, bold=True, color=WHITE)

    for i, row in enumerate(geo_entity[:14]):
        top = Inches(1.95 + i * 0.34)
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        count = int(row.get("pages_with_schema", 0) or 0)
        pct_val = row.get("coverage_pct", "0%")
        pct_num = int(pct_val.replace("%", "") or 0)
        status = "Present" if count > 0 else "Absent"
        status_color = GREEN if count > 0 else SEV_COLORS["Critical"]

        add_rect(slide, Inches(0.35), top, Inches(12.5), Inches(0.3), bg)
        add_text_box(slide, row.get("entity_type", ""),
                     Inches(0.4), top + Inches(0.03), Inches(5.3), Inches(0.28),
                     font_size=9, color=DARK)
        add_text_box(slide, f"{count:,}",
                     Inches(5.8), top + Inches(0.03), Inches(1.7), Inches(0.28),
                     font_size=9, color=DARK)
        add_text_box(slide, pct_val,
                     Inches(7.6), top + Inches(0.03), Inches(1.3), Inches(0.28),
                     font_size=9, color=DARK)
        add_rect(slide, Inches(9.0), top + Inches(0.02), Inches(1.0), Inches(0.26), status_color)
        add_text_box(slide, status, Inches(9.0), top + Inches(0.02), Inches(1.0), Inches(0.26),
                     font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_ai_crawlability(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "AI Crawlability & GEO Signals", brand_rgb,
                     subtitle="Readiness for AI answer engines (ChatGPT, Perplexity, Google AI Overviews)",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    geo_ai = d.get("geo_ai", {})
    signals = [
        ("llms.txt Present", d.get("llms_present", False), True, "GEO — Absent means no AI-readable site summary"),
        ("AI Bots Blocked (robots.txt)", d.get("blocked_ai_bots", "None") not in ("None detected", "Not checked", ""), False, "GEO — Blocked bots cannot crawl for AI citations"),
        ("robots.txt Has Admin Rules", geo_ai.get("robots_has_admin_rule", "N/A") == "True", True, "Technical — Admin paths should be disallowed"),
        ("robots.txt Has JSON:API Rules", geo_ai.get("robots_has_jsonapi_rule", "N/A") == "True", True, "Technical — API endpoints should be disallowed"),
    ]

    for i, (label, is_good_when_true, good_is_true, note) in enumerate(signals):
        top = Inches(1.5 + i * 1.3)
        val_text = "YES" if is_good_when_true else "NO"
        is_good = is_good_when_true == good_is_true
        color = GREEN if is_good else SEV_COLORS["Critical"]
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(1.1), LIGHT_GREY)
        add_rect(slide, Inches(0.4), top, Inches(1.5), Inches(1.1), color)
        add_text_box(slide, val_text, Inches(0.4), top + Inches(0.3),
                     Inches(1.5), Inches(0.5),
                     font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, Inches(2.1), top + Inches(0.1),
                     Inches(10), Inches(0.45),
                     font_size=14, bold=True, color=DARK)
        add_text_box(slide, note, Inches(2.1), top + Inches(0.58),
                     Inches(10), Inches(0.42),
                     font_size=10, color=DARK)


def build_recommendations(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Recommendations — Priority Matrix", brand_rgb,
                     subtitle="Ranked by severity and implementation effort",
                     company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    headers = ["#", "Category", "Issue", "Severity", "Effort", "Impact"]
    col_x = [Inches(0.3), Inches(0.75), Inches(1.8), Inches(8.1), Inches(9.3), Inches(10.4)]
    col_w = [Inches(0.42), Inches(1.0), Inches(6.2), Inches(1.1), Inches(1.0), Inches(1.0)]

    add_rect(slide, Inches(0.3), Inches(1.25), Inches(11.3), Inches(0.38), brand_rgb)
    for h, cx, cw in zip(headers, col_x, col_w):
        add_text_box(slide, h, cx, Inches(1.28), cw, Inches(0.33),
                     font_size=9, bold=True, color=WHITE)

    for i, rec in enumerate(d["recs"][:12]):
        top = Inches(1.65 + i * 0.44)
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.3), top, Inches(11.3), Inches(0.4), bg)

        sev = rec.get("severity", "Low")
        vals = [
            str(rec.get("priority_rank", "")),
            str(rec.get("category", ""))[:10],
            str(rec.get("issue_title", ""))[:75],
            sev,
            str(rec.get("effort", ""))[:8],
            str(rec.get("impact", ""))[:8],
        ]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            if j == 3:
                sev_col = SEV_COLORS.get(sev, DARK)
                add_rect(slide, cx, top + Inches(0.06), cw, Inches(0.28), sev_col)
                add_text_box(slide, val, cx, top + Inches(0.06), cw, Inches(0.28),
                             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            else:
                add_text_box(slide, val, cx + Inches(0.04), top + Inches(0.05),
                             cw - Inches(0.05), Inches(0.35), font_size=9, color=DARK)


def build_next_steps(prs, d, brand_rgb, args, audit_date):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_slide_header(slide, "Next Steps", brand_rgb, company=args.company)
    add_footer(slide, args.company, args.confidential, brand_rgb)

    # Quick wins (< 1 week)
    add_rect(slide, Inches(0.4), Inches(1.25), Inches(6.0), Inches(5.8), LIGHT_GREY)
    add_rect(slide, Inches(0.4), Inches(1.25), Inches(6.0), Inches(0.45), GREEN)
    add_text_box(slide, "Quick Wins (< 1 week)",
                 Inches(0.5), Inches(1.28), Inches(5.8), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)

    quick_wins = [r for r in d["recs"] if r.get("effort") == "Low"][:6]
    for i, rec in enumerate(quick_wins):
        add_text_box(slide, f"• {rec.get('issue_title', '')[:65]}",
                     Inches(0.55), Inches(1.8 + i * 0.8), Inches(5.7), Inches(0.7),
                     font_size=10, color=DARK)

    # Strategic actions (1-3 months)
    add_rect(slide, Inches(6.9), Inches(1.25), Inches(6.0), Inches(5.8), LIGHT_GREY)
    add_rect(slide, Inches(6.9), Inches(1.25), Inches(6.0), Inches(0.45), brand_rgb)
    add_text_box(slide, "Strategic Actions (1–3 months)",
                 Inches(7.0), Inches(1.28), Inches(5.8), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)

    strategic = [r for r in d["recs"] if r.get("effort") in ("Medium", "High")][:6]
    for i, rec in enumerate(strategic):
        add_text_box(slide, f"• {rec.get('issue_title', '')[:65]}",
                     Inches(7.05), Inches(1.8 + i * 0.8), Inches(5.7), Inches(0.7),
                     font_size=10, color=DARK)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    args = build_parser().parse_args()
    input_dir = os.path.expanduser(args.input)
    if not os.path.isdir(input_dir):
        sys.exit(f"Input directory not found: {input_dir}")

    output_path = args.output or os.path.join(input_dir, "seo-audit-report.pptx")
    audit_date = args.date or datetime.now().strftime("%B %Y")
    brand_rgb = hex_to_rgb(args.brand_color)
    args.confidential = getattr(args, "confidential", False)

    print(f"Loading audit data from: {input_dir}")
    d = load_audit_data(input_dir)
    if args.url:
        d["base_url"] = args.url
    print(f"  {d['total']:,} rows in seo_crawl_raw.csv")
    print(f"  {len(d['recs'])} recommendations")
    print(f"  GEO data: {'yes' if d['has_geo_data'] else 'no'}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    build_cover(prs, d, brand_rgb, args, audit_date)
    build_agenda(prs, d, brand_rgb, args, audit_date)
    build_methodology(prs, d, brand_rgb, args, audit_date)   # Always included
    build_scope_summary(prs, d, brand_rgb, args, audit_date)
    build_seo_dashboard(prs, d, brand_rgb, args, audit_date)
    build_critical_issues(prs, d, brand_rgb, args, audit_date)
    build_onpage_seo(prs, d, brand_rgb, args, audit_date)
    build_technical_seo(prs, d, brand_rgb, args, audit_date)

    if d["has_geo_data"] and not args.no_geo:
        build_geo_slide(prs, d, brand_rgb, args, audit_date)
        build_ai_crawlability(prs, d, brand_rgb, args, audit_date)

    build_recommendations(prs, d, brand_rgb, args, audit_date)
    build_next_steps(prs, d, brand_rgb, args, audit_date)

    prs.save(output_path)
    slide_count = len(prs.slides)
    print(f"\n=== PowerPoint Generated ===")
    print(f"  Slides    : {slide_count}")
    print(f"  Output    : {output_path}")
    print(f"  File size : {os.path.getsize(output_path):,} bytes")


if __name__ == "__main__":
    main()
